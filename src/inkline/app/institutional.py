"""Institutional fixture helpers for editable PPTX sign-off workflows.

This module adds a small execution surface for the weekend vertical slice:

- load YAML / JSON slide specs
- render PDF and PPTX artifacts from those specs
- inspect PPTX editability / fallback heuristics
- render PPTX to PDF via soffice
- audit rendered PPTX artifacts
- compare rendered PDF outputs for parity
"""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any


DEFAULT_OUTPUT_ROOT = Path("~/.local/share/inkline/output").expanduser()
FIXTURE_PAGE_MAP = {
    "cover": 1,
    "team_grid": 4,
    "institutional_kpi_cards": 5,
    "institutional_timeline": 6,
    "appendix_matrix": 7,
}


@dataclass
class RenderArtifacts:
    pdf_path: Path | None = None
    pptx_path: Path | None = None
    export_metadata_path: Path | None = None


def _portable_sidecar_payload(value: Any) -> Any:
    """Strip host-bound absolute-path fields from exported sidecar metadata."""
    if isinstance(value, dict):
        cleaned: dict[str, Any] = {}
        for key, item in value.items():
            if key in {"source_name", "deck_ref", "storyboard_path", "authoring_trace_path"}:
                continue
            cleaned[key] = _portable_sidecar_payload(item)
        return cleaned
    if isinstance(value, list):
        return [_portable_sidecar_payload(item) for item in value]
    return value


def load_spec_file(path: str | Path) -> dict[str, Any]:
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(path)
    suffix = path.suffix.lower()
    text = path.read_text(encoding="utf-8")
    if suffix in {".yaml", ".yml"}:
        try:
            import yaml
        except ImportError as exc:
            raise RuntimeError(
                "YAML spec support requires PyYAML. Install it or use JSON fixtures instead."
            ) from exc
        data = yaml.safe_load(text) or {}
    elif suffix == ".json":
        data = json.loads(text)
    else:
        raise ValueError(f"Unsupported spec file type: {path.suffix}")
    if not isinstance(data, dict):
        raise ValueError("Spec root must be a mapping")
    return load_spec_file_dict(data, source_name=str(path))


def render_spec_file(
    spec_path: str | Path,
    *,
    formats: list[str],
    output_dir: str | Path | None = None,
    editable_institutional: bool = False,
    brand_override: str = "",
    template_override: str = "",
    execution_mode: str = "",
    design_locked: bool | None = None,
    use_design_advisor: bool | None = None,
    authoring_mode: str = "",
) -> RenderArtifacts:
    from inkline.pptx import export_pptx_slides
    from inkline.typst import export_typst_slides

    spec = load_spec_file(spec_path)
    if execution_mode:
        spec["execution_mode"] = execution_mode
    if design_locked is not None:
        spec["design_locked"] = design_locked
    if use_design_advisor is not None:
        spec["use_design_advisor"] = use_design_advisor
    if authoring_mode:
        spec["authoring_mode"] = authoring_mode
    spec = load_spec_file_dict(spec, source_name=str(spec_path))
    slides = spec.get("slides", [])
    title = str(spec.get("title", Path(spec_path).stem))
    brand = str(brand_override or spec.get("brand", "minimal"))
    template = str(template_override or spec.get("template", "consulting"))

    out_dir = Path(output_dir) if output_dir else DEFAULT_OUTPUT_ROOT / Path(spec_path).stem
    out_dir.mkdir(parents=True, exist_ok=True)
    stem = Path(spec_path).stem

    artifacts = RenderArtifacts()
    artifact_paths = {}
    try:
        from inkline.intelligence.storyboard import write_storyboard_artifacts
        artifact_paths = write_storyboard_artifacts(spec, output_dir=out_dir, stem=stem)
    except Exception as exc:
        artifact_paths = {}
        print(
            f"[inkline render] WARNING: could not write storyboard artifacts: {exc}",
            file=sys.stderr,
        )
        raise

    if "pdf" in formats:
        artifacts.pdf_path = out_dir / f"{stem}.pdf"
        export_typst_slides(
            slides=slides,
            output_path=str(artifacts.pdf_path),
            brand=brand,
            template=template,
        )

    if "pptx" in formats:
        artifacts.pptx_path = out_dir / f"{stem}.pptx"
        artifacts.export_metadata_path = out_dir / f"{stem}.export_metadata.json"
        export_pptx_slides(
            slides=slides,
            output_path=artifacts.pptx_path,
            brand=brand,
            title=title,
            source_root=Path(spec_path).parent,
            metadata_path=artifacts.export_metadata_path,
            editable_institutional=editable_institutional,
            deck_metadata=_portable_sidecar_payload({
                "storyboard": spec.get("_resolved_storyboard", {}),
                "authoring_trace": spec.get("_authoring_trace", {}),
                # Keep the export sidecar portable across machines.
                "artifact_files": {
                    "storyboard": artifact_paths.get("storyboard_path", Path("")).name
                    if artifact_paths.get("storyboard_path")
                    else "",
                    "authoring_trace": artifact_paths.get("authoring_trace_path", Path("")).name
                    if artifact_paths.get("authoring_trace_path")
                    else "",
                },
            }),
        )

    return artifacts


def load_spec_file_dict(spec: dict[str, Any], *, source_name: str = "") -> dict[str, Any]:
    from inkline.intelligence.storyboard import resolve_storyboard_spec

    data = dict(spec)
    data.setdefault("slides", [])
    return resolve_storyboard_spec(data, source_name=source_name)


def inspect_pptx(pptx_path: str | Path) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    export_meta_path = pptx_path.with_suffix(".export_metadata.json")
    export_meta = None
    if export_meta_path.exists():
        try:
            export_meta = json.loads(export_meta_path.read_text(encoding="utf-8"))
        except Exception:
            export_meta = None
    slide_entries: list[dict[str, Any]] = []
    editable_count = 0
    fully_native_count = 0
    fallback_count = 0
    exception_count = 0

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for idx, slide in enumerate(prs.slides, start=1):
        picture_like = 0
        full_slide_picture = False
        shape_count = len(slide.shapes)
        text_shapes = 0

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text_shapes += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_like += 1
                if (
                    abs(shape.left) < 1000
                    and abs(shape.top) < 1000
                    and abs(shape.width - slide_width) < 2000
                    and abs(shape.height - slide_height) < 2000
                ):
                    full_slide_picture = True

        status = "native"
        fallback_reason = ""
        editability_exceptions: list[str] = []
        if export_meta and idx <= len(export_meta.get("slides", [])):
            exported = export_meta["slides"][idx - 1]
            status = exported.get("status", status)
            fallback_reason = exported.get("fallback_reason", "")
            editability_exceptions = list(exported.get("pptx_editability_exceptions", []))
            if status == "fallback":
                fallback_count += 1
            elif status == "native_with_exceptions":
                exception_count += 1
                editable_count += 1
            else:
                editable_count += 1
                fully_native_count += 1
        elif full_slide_picture and shape_count <= 2:
            status = "fallback"
            fallback_reason = "full_slide_picture"
            fallback_count += 1
        else:
            editable_count += 1
            fully_native_count += 1

        slide_entries.append(
            {
                "slide_number": idx,
                "shape_count": shape_count,
                "text_shape_count": text_shapes,
                "picture_count": picture_like,
                "status": status,
                "fallback_reason": fallback_reason,
                "pptx_editability_exceptions": editability_exceptions,
            }
        )

    total = len(slide_entries) or 1
    inspection_mode = "metadata" if export_meta else "heuristic"
    return {
        "pptx_path": str(pptx_path),
        "slide_count": len(slide_entries),
        "editable_native_ratio": editable_count / total,
        "fully_native_ratio": fully_native_count / total,
        "slides_with_image_fallback": [
            s["slide_number"] for s in slide_entries if s["status"] == "fallback"
        ],
        "slides_with_editability_exceptions": [
            s["slide_number"] for s in slide_entries if s["status"] == "native_with_exceptions"
        ],
        "inspection_mode": inspection_mode,
        "reliability": "high" if export_meta else "best_effort",
        "warning": (
            ""
            if export_meta
            else "Best-effort PPTX inspection without adjacent export metadata; "
            "native/fallback classification may be inaccurate."
        ),
        "fallback_reasons": {
            str(s["slide_number"]): s["fallback_reason"]
            for s in slide_entries
            if s["fallback_reason"]
        },
        "slide_statuses": {
            str(s["slide_number"]): s["status"] for s in slide_entries
        },
        "slides": slide_entries,
    }


def render_pptx_via_soffice(
    pptx_path: str | Path,
    *,
    output_dir: str | Path | None = None,
    timeout: int = 90,
) -> Path:
    pptx_path = Path(pptx_path)
    soffice = shutil.which("soffice")
    if not soffice:
        raise RuntimeError("soffice not found on PATH")
    out_dir = Path(output_dir) if output_dir else Path(tempfile.mkdtemp(prefix="inkline_soffice_"))
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(pptx_path),
    ]
    proc = subprocess.run(
        cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=timeout,
        check=False,
    )
    soffice_pdf = out_dir / f"{pptx_path.stem}.pdf"
    if proc.returncode != 0 or not soffice_pdf.exists():
        raise RuntimeError(
            f"soffice conversion failed rc={proc.returncode}: "
            f"{proc.stdout[-300:]} {proc.stderr[-300:]}"
        )
    rendered = out_dir / f"{pptx_path.stem}.rendered.pdf"
    if rendered.exists():
        rendered.unlink()
    soffice_pdf.rename(rendered)
    return rendered


def audit_pptx(
    pptx_path: str | Path,
    *,
    rubric: str = "institutional",
    brand: str = "",
    output_path: str | Path | None = None,
    rendered_pdf_path: str | Path | None = None,
) -> dict[str, Any]:
    from inkline.intelligence.vishwakarma import critique_pdf

    pptx_path = Path(pptx_path)
    if rendered_pdf_path:
        rendered_pdf = Path(rendered_pdf_path)
    else:
        rendered_tmp = render_pptx_via_soffice(pptx_path)
        rendered_pdf = pptx_path.with_suffix(".rendered.pdf")
        rendered_pdf.write_bytes(rendered_tmp.read_bytes())

    provider = os.environ.get("INKLINE_VISION_PROVIDER", "")
    if not provider:
        provider = "codex_cli"
        os.environ["INKLINE_VISION_PROVIDER"] = provider
    result = critique_pdf(str(rendered_pdf), rubric=rubric, brand=brand).to_dict()
    result["artifact_type"] = "pptx_render"
    result["pptx_path"] = str(pptx_path)
    result["rendered_pdf_path"] = str(rendered_pdf)
    result["provider_trace"] = [{"provider": provider, "source": "env/requested"}]

    export_meta_path = pptx_path.with_suffix(".export_metadata.json")
    if export_meta_path.exists():
        try:
            export_meta = json.loads(export_meta_path.read_text(encoding="utf-8"))
            result["export_metadata_path"] = str(export_meta_path)
            result["storyboard_schema_version"] = (
                export_meta.get("deck_metadata", {})
                .get("storyboard", {})
                .get("schema_version")
            )
            result["storyboard_audit"] = _aggregate_storyboard_audit(
                result=result,
                slide_meta=export_meta.get("slides", []),
                fallback_key="status",
            )
        except Exception as exc:
            result["storyboard_audit_error"] = str(exc)
    else:
        result["storyboard_audit"] = {
            "schema_name": "deck_audit",
            "schema_version": 1,
            "deck_verdict": "needs_human_signoff",
            "deck_required_fix_count": 1,
            "slides_failed_hard_checks": [],
            "slides_requiring_human_signoff": [],
            "dimensions_not_evaluated": [
                {
                    "slide_index": None,
                    "dimensions": ["visual_quality_metadata_context", "storyboard_metadata"],
                }
            ],
            "warning_budget_used": 0,
            "slides": [],
            "reason": "PPTX export metadata sidecar missing; storyboard-aware audit dimensions unavailable.",
        }

    if output_path:
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(result, indent=2), encoding="utf-8")
    return result


def audit_pdf_artifact(
    pdf_path: str | Path,
    *,
    rubric: str = "institutional",
    brand: str = "",
    output_path: str | Path | None = None,
) -> dict[str, Any]:
    from inkline.intelligence.vishwakarma import critique_pdf

    pdf_path = Path(pdf_path)
    provider = os.environ.get("INKLINE_VISION_PROVIDER", "")
    if not provider:
        provider = "codex_cli"
        os.environ["INKLINE_VISION_PROVIDER"] = provider
    result = critique_pdf(str(pdf_path), rubric=rubric, brand=brand).to_dict()
    result["artifact_type"] = "pdf"
    result["pdf_path"] = str(pdf_path)
    result["provider_trace"] = [{"provider": provider, "source": "env/requested"}]

    storyboard_path = pdf_path.with_name(f"{pdf_path.stem}.storyboard.json")
    if storyboard_path.exists():
        try:
            storyboard = json.loads(storyboard_path.read_text(encoding="utf-8"))
            result["storyboard_schema_version"] = storyboard.get("schema_version")
            result["storyboard_audit"] = _aggregate_storyboard_audit(
                result=result,
                slide_meta=storyboard.get("slides", []),
                fallback_key="fallback_used",
            )
        except Exception as exc:
            result["storyboard_audit_error"] = str(exc)
    else:
        result["storyboard_audit"] = {
            "schema_name": "deck_audit",
            "schema_version": 1,
            "deck_verdict": "needs_human_signoff",
            "deck_required_fix_count": 1,
            "slides_failed_hard_checks": [],
            "slides_requiring_human_signoff": [],
            "dimensions_not_evaluated": [
                {
                    "slide_index": None,
                    "dimensions": ["storyboard_metadata"],
                }
            ],
            "warning_budget_used": 0,
            "slides": [],
            "reason": "Storyboard metadata missing; storyboard-aware audit dimensions unavailable.",
        }

    if output_path:
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(result, indent=2), encoding="utf-8")
    return result


def _aggregate_storyboard_audit(
    *,
    result: dict[str, Any],
    slide_meta: list[dict[str, Any]],
    fallback_key: str,
) -> dict[str, Any]:
    from inkline.intelligence.audit_storyboard import aggregate_deck_audit, evaluate_slide_audit

    critiques_by_index = {
        int(critique.get("slide_index", 0)): critique
        for critique in result.get("slide_critiques", [])
        if int(critique.get("slide_index", 0)) > 0
    }
    slide_results = []
    total_slides = max(len(slide_meta), max(critiques_by_index, default=0))
    for idx in range(1, total_slides + 1):
        critique = critiques_by_index.get(idx, {})
        meta = slide_meta[idx - 1] if idx <= len(slide_meta) else {}
        storyboard = meta.get("storyboard", meta)
        critique_verdict = (
            "INCOMPLETE"
            if not critique or not meta
            else str(critique.get("verdict", "INCOMPLETE"))
        )
        slide_results.append(
            evaluate_slide_audit(
                slide_index=idx,
                storyboard=storyboard,
                critique_verdict=critique_verdict,
                archetype_declared=bool((storyboard or {}).get("archetype")),
                reference_family_declared=bool((storyboard or {}).get("reference_family")),
                fallback_used=bool(meta.get(fallback_key) == "fallback" if fallback_key == "status" else meta.get(fallback_key)),
            )
        )
    return aggregate_deck_audit(slide_results)


def compare_rendered_pdfs(
    baseline_pdf: str | Path,
    rendered_pdf: str | Path,
    *,
    slide_tokens: list[str],
    output_path: str | Path | None = None,
) -> dict[str, Any]:
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        raise RuntimeError("pdftoppm not found on PATH")

    with tempfile.TemporaryDirectory(prefix="inkline_compare_") as tmp:
        tmpdir = Path(tmp)
        base_prefix = tmpdir / "base"
        render_prefix = tmpdir / "render"
        subprocess.run([pdftoppm, "-r", "144", "-png", str(baseline_pdf), str(base_prefix)], check=True)
        subprocess.run([pdftoppm, "-r", "144", "-png", str(rendered_pdf), str(render_prefix)], check=True)

        base_images = sorted(tmpdir.glob("base-*.png"))
        render_images = sorted(tmpdir.glob("render-*.png"))

        page_indices = _resolve_slide_tokens(slide_tokens, base_images, render_images)
        slide_scores: list[dict[str, Any]] = []
        for label, page_index in page_indices:
            score = _compare_images(base_images[page_index - 1], render_images[page_index - 1])
            slide_scores.append({"slide": label, "page_index": page_index, "score": score})

    parity_score = sum(item["score"] for item in slide_scores) / max(len(slide_scores), 1)
    report = {
        "baseline_pdf": str(baseline_pdf),
        "rendered_pdf": str(rendered_pdf),
        "parity_diff_score": parity_score,
        "slides": slide_scores,
    }
    if output_path:
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(report, indent=2), encoding="utf-8")
    return report


def _resolve_slide_tokens(
    tokens: list[str],
    base_images: list[Path],
    render_images: list[Path],
) -> list[tuple[str, int]]:
    if len(base_images) != len(render_images):
        raise RuntimeError("baseline and rendered PDFs produced different page counts")
    resolved: list[tuple[str, int]] = []
    for token in tokens:
        raw = token.strip()
        if not raw:
            continue
        if raw.isdigit():
            idx = int(raw)
        else:
            idx = FIXTURE_PAGE_MAP.get(raw)
            if idx is None:
                raise RuntimeError(f"Unknown slide token for parity compare: {raw}")
        if idx < 1 or idx > len(base_images):
            raise RuntimeError(f"Slide/page index out of range for parity compare: {raw} -> {idx}")
        resolved.append((raw, idx))
    return resolved


def _compare_images(lhs: Path, rhs: Path) -> float:
    from PIL import Image, ImageChops

    a = Image.open(lhs).convert("RGB")
    b = Image.open(rhs).convert("RGB")
    if a.size != b.size:
        b = b.resize(a.size)
    diff = ImageChops.difference(a, b)
    hist = diff.histogram()
    total = 0
    max_diff = a.size[0] * a.size[1] * 255 * 3
    for value, count in enumerate(hist):
        channel_value = value % 256
        total += channel_value * count
    return total / max_diff


def dump_json(path: str | Path, payload: dict[str, Any]) -> Path:
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    return path
