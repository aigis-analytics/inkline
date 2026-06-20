"""PPTX-first reference deck ingestion for local/private benchmark catalogs."""

from __future__ import annotations

import base64
import json
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from inkline.intelligence.reference_catalog import LOCAL_ROOT, validate_reference_family_id
from inkline.intelligence.full_slide_archetypes import (
    KNOWN_SLIDE_ROLES,
    ROLE_TO_DEFAULT_ARCHETYPE,
    get_full_slide_archetype,
)
from inkline.intelligence.reference_schema import validate_reference_slide_manifest

PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVQIHWP4////fwAJ+wP9KobjigAAAABJRU5ErkJggg=="
)


def _write_placeholder_png(path: Path) -> None:
    path.write_bytes(PNG_1X1)


def _parse_simple_yaml(text: str) -> dict[str, Any]:
    result: dict[str, Any] = {}
    current_list: list[dict[str, Any]] | None = None
    current_item: dict[str, Any] | None = None
    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            continue
        if line.startswith("reference_family_id:"):
            result["reference_family_id"] = line.split(":", 1)[1].strip()
        elif line.strip() == "slides:":
            current_list = []
            result["slides"] = current_list
        elif line.lstrip().startswith("- ") and current_list is not None:
            current_item = {}
            current_list.append(current_item)
            key, _, value = line.lstrip()[2:].partition(":")
            current_item[key.strip()] = value.strip()
        elif current_item is not None and ":" in line:
            key, _, value = line.strip().partition(":")
            current_item[key.strip()] = value.strip()
    return result


def _safe_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()
    return ""


def _normalize_shape(shape: Any, slide_width: int, slide_height: int) -> dict[str, Any]:
    return {
        "type": getattr(shape, "shape_type", None).name if getattr(shape, "shape_type", None) else "UNKNOWN",
        "x": round(float(getattr(shape, "left", 0)) / float(slide_width or 1), 6),
        "y": round(float(getattr(shape, "top", 0)) / float(slide_height or 1), 6),
        "w": round(float(getattr(shape, "width", 0)) / float(slide_width or 1), 6),
        "h": round(float(getattr(shape, "height", 0)) / float(slide_height or 1), 6),
        "text": _safe_text(shape),
    }


def _infer_composition_family(
    *,
    normalized: list[dict[str, Any]],
    text_blocks: list[dict[str, Any]],
    role: str,
) -> str:
    picture_count = sum(1 for item in normalized if "PICTURE" in str(item.get("type", "")).upper())
    width_heavy = max((float(item.get("w", 0.0) or 0.0) for item in normalized), default=0.0)
    if role == "cover" and width_heavy >= 0.45:
        return "hero_cover"
    if role in {"team", "people", "key_people"} or picture_count >= 3:
        return "people_profiles"
    if role in {"timeline", "process", "roadshow"}:
        return "timeline_spine"
    if role in {"appendix_ranked_table", "pipeline"} or len(text_blocks) >= 8:
        return "dense_table"
    if role in {"economics", "size_of_prize"}:
        return "two_zone_summary"
    return "card_grid"


def _infer_density_class(*, normalized: list[dict[str, Any]], text_blocks: list[dict[str, Any]]) -> str:
    score = len(normalized) + len(text_blocks)
    if score >= 16:
        return "high"
    if score >= 7:
        return "medium"
    return "low"


def _infer_zone_map(normalized: list[dict[str, Any]]) -> dict[str, Any]:
    if not normalized:
        return {}
    title_like = next((item for item in normalized if item.get("text")), None)
    title_zone = {
        "x": title_like.get("x", 0.0),
        "y": title_like.get("y", 0.0),
        "w": title_like.get("w", 0.0),
        "h": title_like.get("h", 0.0),
    } if title_like else {}
    return {
        "title_zone": title_zone,
        "header_present": bool(title_like and float(title_like.get("y", 1.0)) < 0.18),
        "footer_present": any(float(item.get("y", 0.0)) > 0.85 for item in normalized),
        "dominant_axis": "horizontal" if sum(float(item.get("w", 0.0)) for item in normalized) >= sum(float(item.get("h", 0.0)) for item in normalized) else "vertical",
    }


def _infer_content_slots(*, normalized: list[dict[str, Any]], text_blocks: list[dict[str, Any]]) -> list[dict[str, Any]]:
    picture_count = sum(1 for item in normalized if "PICTURE" in str(item.get("type", "")).upper())
    slots = [
        {"name": "text_blocks", "kind": "text", "required": False, "count": len(text_blocks)},
        {"name": "images", "kind": "image", "required": False, "count": picture_count},
    ]
    if any(float(item.get("w", 0.0)) > 0.4 for item in normalized):
        slots.append({"name": "hero_zone", "kind": "hero", "required": False, "count": 1})
    return slots


def _infer_style_tokens(
    *,
    normalized: list[dict[str, Any]],
    role: str,
    prs_style_tokens: dict[str, Any],
) -> dict[str, Any]:
    tokens = dict(prs_style_tokens)
    picture_count = sum(1 for item in normalized if "PICTURE" in str(item.get("type", "")).upper())
    tokens.update(
        {
            "title_alignment": "left",
            "card_count": max(0, len(normalized) - picture_count),
            "picture_count": picture_count,
            "hero_occupancy_ratio": round(max((float(item.get("w", 0.0)) * float(item.get("h", 0.0)) for item in normalized), default=0.0), 4),
            "typography_treatment_class": "editorial_cover" if role == "cover" else "institutional_grid",
            "image_treatment": "headshots" if role in {"team", "people"} else "reference_exhibit" if role in {"market_map", "map"} else "mixed",
        }
    )
    return tokens


def _infer_reference_role(
    *,
    slide_index: int,
    text_blocks: list[dict[str, Any]],
    shape_count: int,
) -> str:
    text = " ".join(block.get("text", "") for block in text_blocks).lower()
    if slide_index == 1:
        return "cover"
    if any(token in text for token in ("team", "boots on the ground", "leadership", "management", "people")):
        return "team"
    if any(token in text for token in ("timeline", "roadmap", "milestone")):
        return "timeline"
    if any(token in text for token in ("process", "execution", "workstream")):
        return "process"
    if any(token in text for token in ("appendix", "pipeline", "matrix", "ranking")):
        return "appendix_ranked_table"
    if shape_count >= 8 and len(text_blocks) >= 4:
        return "appendix_ranked_table"
    if len(text_blocks) <= 2:
        return "cover"
    return "content"


def _style_tokens(prs: Any) -> dict[str, Any]:
    width = int(getattr(prs, "slide_width", 0) or 0)
    height = int(getattr(prs, "slide_height", 0) or 0)
    aspect_ratio = round((width / height), 4) if height else 0
    return {
        "source_mode": "pptx_native_mvp",
        "preview_mode": "placeholder_png",
        "geometry_mode": "normalized_native",
        "aspect_ratio": aspect_ratio,
        "slide_size_emu": {"width": width, "height": height},
    }


def _redacted_source_label(pptx_path: Path) -> str:
    return pptx_path.name


def ingest_reference_pptx(
    pptx_path: str | Path,
    *,
    family_id: str,
    license_classification: str = "private_internal",
    notes: list[str] | None = None,
    catalog_root: str | Path | None = None,
) -> dict[str, Any]:
    """Ingest a PPTX into the local reference catalog.

    The MVP is intentionally PPTX-only. Confidential and private assets are
    always written into the local catalog, never packaged paths.
    """
    pptx_path = Path(pptx_path)
    if pptx_path.suffix.lower() != ".pptx":
        raise ValueError("MVP reference ingestion supports PPTX only")
    if license_classification not in {"public_reusable", "public_reference_only", "private_internal", "client_confidential"}:
        raise ValueError(f"Unsupported license classification: {license_classification}")
    family_id = validate_reference_family_id(family_id)
    if catalog_root is not None and Path(catalog_root).expanduser().resolve() != LOCAL_ROOT.resolve():
        raise ValueError("Reference ingestion may only write to the local private catalog root")
    target_root = LOCAL_ROOT
    target_family = target_root / family_id
    target_family.mkdir(parents=True, exist_ok=True)

    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide_entries: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        reference_slide_id = f"{family_id}_s{idx:02d}"
        preview_path = target_family / f"{reference_slide_id}.png"
        _write_placeholder_png(preview_path)
        normalized = [_normalize_shape(shape, slide_width, slide_height) for shape in slide.shapes]
        text_blocks = [entry for entry in normalized if entry.get("text")]
        inferred_role = _infer_reference_role(
            slide_index=idx,
            text_blocks=text_blocks,
            shape_count=len(normalized),
        )
        composition_family = _infer_composition_family(
            normalized=normalized,
            text_blocks=text_blocks,
            role=inferred_role,
        )
        density_class = _infer_density_class(normalized=normalized, text_blocks=text_blocks)
        manifest = validate_reference_slide_manifest(
            {
                "reference_slide_id": reference_slide_id,
                "reference_family_id": family_id,
                "source_slide_index": idx,
                "role": inferred_role,
                "composition_family": composition_family,
                "density_class": density_class,
                "style_tokens": _infer_style_tokens(
                    normalized=normalized,
                    role=inferred_role,
                    prs_style_tokens=_style_tokens(prs),
                ),
                "zone_map": _infer_zone_map(normalized),
                "content_slots": _infer_content_slots(normalized=normalized, text_blocks=text_blocks),
                "usable_for_retrieval": True,
                "archetype_tag": ROLE_TO_DEFAULT_ARCHETYPE.get(inferred_role, ""),
                "hero_kind": "portraits" if inferred_role in {"team", "people"} else "full_bleed" if inferred_role == "cover" else "",
                "evidence_kind": "table" if density_class == "high" else "cards",
                "benchmark_quality_weight": 0.7 if inferred_role != "content" else 0.45,
                "strong_exemplar": inferred_role != "content",
                "do_not_imitate": False,
                "preview_path": f"{reference_slide_id}.png",
                "curation_notes": [],
                "normalized_geometry": normalized,
                "text_blocks": text_blocks,
            }
        )
        (target_family / f"{reference_slide_id}.json").write_text(
            json.dumps(manifest, indent=2, sort_keys=True), encoding="utf-8"
        )
        slide_entries.append(
            {
                "reference_slide_id": reference_slide_id,
                "role": inferred_role,
                "archetype_candidate": ROLE_TO_DEFAULT_ARCHETYPE.get(
                    inferred_role, ""
                ),
                "composition_family": composition_family,
                "density_class": density_class,
                "benchmark_quality_weight": manifest["benchmark_quality_weight"],
                "strong_exemplar": manifest["strong_exemplar"],
                "do_not_imitate": manifest["do_not_imitate"],
                "usable_for_retrieval": True,
                "preview_path": f"{reference_slide_id}.png",
                "manifest_path": f"{reference_slide_id}.json",
            }
        )

    family_manifest = {
        "schema_name": "reference_family_manifest",
        "schema_version": 1,
        "reference_family_id": family_id,
        "source_id": pptx_path.stem,
        "source_path": _redacted_source_label(pptx_path),
        "license_classification": license_classification,
        "ingestion_method": "pptx_native",
        "confidence_score": 1.0,
        "ingested_at": datetime.now(UTC).isoformat(),
        "version": 1,
        "curated_by": "",
        "style_tokens": _style_tokens(prs),
        "slides": slide_entries,
        "notes": list(notes or []),
        "override_log": [],
    }
    manifest_path = target_family / "reference_family_manifest.json"
    manifest_path.write_text(json.dumps(family_manifest, indent=2, sort_keys=True), encoding="utf-8")
    overrides_path = target_family / "curation_overrides.yaml"
    if not overrides_path.exists():
        overrides_path.write_text(
            "reference_family_id: %s\nslides: []\n" % family_id,
            encoding="utf-8",
        )
    return family_manifest


def apply_curation_overrides(reference_family_id: str, *, catalog_root: str | Path | None = None) -> dict[str, Any]:
    reference_family_id = validate_reference_family_id(reference_family_id)
    if catalog_root is not None and Path(catalog_root).expanduser().resolve() != LOCAL_ROOT.resolve():
        raise ValueError("Reference curation may only read/write from the local private catalog root")
    target_root = LOCAL_ROOT
    family_dir = target_root / reference_family_id
    manifest_path = family_dir / "reference_family_manifest.json"
    overrides_path = family_dir / "curation_overrides.yaml"
    if not manifest_path.exists():
        raise FileNotFoundError(manifest_path)
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    if not overrides_path.exists():
        return payload
    try:
        import yaml
        overrides = yaml.safe_load(overrides_path.read_text(encoding="utf-8")) or {}
    except ImportError:
        overrides = _parse_simple_yaml(overrides_path.read_text(encoding="utf-8"))
    slides_by_id = {slide["reference_slide_id"]: slide for slide in payload.get("slides", [])}
    for item in overrides.get("slides", []) or []:
        slide = slides_by_id.get(item.get("reference_slide_id"))
        if not slide:
            continue
        role_override = item.get("role_override")
        if role_override and str(role_override) not in KNOWN_SLIDE_ROLES:
            raise ValueError(f"Unknown role_override '{role_override}' for {slide['reference_slide_id']}")
        archetype_override = item.get("archetype_override")
        if archetype_override:
            get_full_slide_archetype(str(archetype_override))
        for source, target in (
            ("role_override", "role"),
            ("archetype_override", "archetype_candidate"),
            ("composition_family", "composition_family"),
            ("density_class", "density_class"),
            ("benchmark_quality_weight", "benchmark_quality_weight"),
            ("strong_exemplar", "strong_exemplar"),
            ("do_not_imitate", "do_not_imitate"),
            ("usable_for_retrieval", "usable_for_retrieval"),
            ("exemplar_strength", "exemplar_strength"),
            ("imitate", "imitate"),
            ("notes", "notes"),
        ):
            if source in item:
                slide[target] = item[source]
        payload.setdefault("override_log", []).append(
            {
                "reference_slide_id": slide["reference_slide_id"],
                "applied_fields": sorted(k for k in item.keys() if k != "reference_slide_id"),
            }
        )
    manifest_path.write_text(json.dumps(payload, indent=2, sort_keys=True), encoding="utf-8")
    return payload


__all__ = ["apply_curation_overrides", "ingest_reference_pptx"]
